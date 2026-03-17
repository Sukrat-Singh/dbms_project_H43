"""
generate_ppt.py  –  Generate a 10-slide PowerPoint presentation for the
Lab Test Management System DBMS project.

Run:  python generate_ppt.py
Output: Lab_Test_Management_Presentation.pptx (in the same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ══════════════════════════════════════════════════════════════════════
#  COLOR PALETTE
# ══════════════════════════════════════════════════════════════════════
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x96, 0xD6)   # Vivid blue
ACCENT_GREEN  = RGBColor(0x00, 0xC9, 0xA7)   # Teal green
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)   # Warm orange
ACCENT_PURPLE = RGBColor(0x7B, 0x68, 0xEE)   # Medium purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT      = RGBColor(0x2D, 0x2D, 0x2D)
SECTION_BG     = RGBColor(0x0F, 0x0F, 0x1A)  # Darker bg for variety
CARD_BG        = RGBColor(0x25, 0x25, 0x3A)   # Card background


def add_bg(slide, color=DARK_BG):
    """Fill the entire slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with a single formatted paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items,
                       font_size=16, color=WHITE, font_name="Calibri",
                       bullet_color=ACCENT_BLUE, spacing=Pt(8)):
    """Add a text box with multiple bullet-point paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox


def add_table_slide(slide, left, top, width, height, headers, rows,
                    header_color=ACCENT_BLUE, row_color1=CARD_BG,
                    row_color2=DARK_BG):
    """Add a formatted table to a slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # Style data rows
    for i, row in enumerate(rows):
        bg = row_color1 if i % 2 == 0 else row_color2
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = WHITE
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    return table_shape


def add_slide_number(slide, num, total=10):
    """Add a subtle slide number at the bottom-right."""
    add_textbox(slide,
                Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num} / {total}", font_size=10, color=LIGHT_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_accent_bar(slide, top=Inches(1.55), color=ACCENT_BLUE):
    """Add a thin accent bar under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(2), Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ══════════════════════════════════════════════════════════════════════
#  BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════
def build_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 1: TITLE
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, SECTION_BG)

    # Decorative accent circles
    for cx, cy, sz, clr in [
        (8.2, 0.3, 1.8, ACCENT_PURPLE), (8.8, 1.5, 1.0, ACCENT_BLUE),
        (0.2, 5.5, 1.5, ACCENT_GREEN),  (1.2, 6.2, 0.8, ACCENT_BLUE),
    ]:
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(sz), Inches(sz)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = clr
        circ.fill.fore_color.brightness = 0.6
        circ.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(1),
                "🧪", font_size=52, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(8), Inches(1.2),
                "Lab Test Management System", font_size=40, bold=True,
                color=WHITE, font_name="Calibri")
    add_textbox(slide, Inches(0.8), Inches(3.4), Inches(8), Inches(0.6),
                "A MongoDB-Based DBMS Project", font_size=22,
                color=ACCENT_BLUE, font_name="Calibri")

    # Subtle info line
    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.5),
                "Course: Database Management Systems   |   Group: H43",
                font_size=14, color=LIGHT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(8), Inches(0.5),
                "Technologies: Python  •  MongoDB Atlas  •  Streamlit",
                font_size=14, color=LIGHT_GRAY)
    add_slide_number(slide, 1)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 2: PROBLEM STATEMENT & INTRODUCTION
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Problem Statement & Introduction", font_size=32, bold=True,
                color=ACCENT_BLUE)
    add_accent_bar(slide, Inches(1.2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.2), Inches(1.2),
                "Hospitals and diagnostic labs handle thousands of lab tests daily. "
                "Managing patient records, test catalogs, orders, results, and alerts "
                "manually or with legacy systems leads to errors, delays, and missed "
                "critical results. Our system digitizes this entire workflow.",
                font_size=15, color=LIGHT_GRAY)

    # Feature cards
    features = [
        ("👤 Patient\nManagement", "Register, search,\nand manage patients", ACCENT_BLUE),
        ("🔬 Test\nCatalog", "Define tests with\nnormal ranges & costs", ACCENT_GREEN),
        ("📋 Order\nTracking", "Place orders, track\nstatus in real-time", ACCENT_ORANGE),
        ("🔔 Smart\nAlerts", "Auto-alerts on critical\nor delayed results", ACCENT_PURPLE),
    ]
    for i, (title, desc, color) in enumerate(features):
        left = Inches(0.8 + i * 2.25)
        card = add_shape_rect(slide, left, Inches(3.2), Inches(2.0), Inches(2.5),
                              CARD_BG, color)

        add_textbox(slide, left + Inches(0.15), Inches(3.4), Inches(1.7), Inches(1.0),
                    title, font_size=15, bold=True, color=color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), Inches(4.4), Inches(1.7), Inches(1.0),
                    desc, font_size=12, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(8), Inches(0.5),
                "📊 Plus: Revenue analytics, turnaround monitoring, and a Streamlit dashboard",
                font_size=13, color=ACCENT_GREEN)
    add_slide_number(slide, 2)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 3: TECHNOLOGY STACK
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Technology Stack", font_size=32, bold=True, color=ACCENT_GREEN)
    add_accent_bar(slide, Inches(1.2), ACCENT_GREEN)

    techs = [
        ("🐍  Python 3.x",     "Core programming language for all backend logic,\n"
                                "services, triggers, and stored procedures"),
        ("🍃  MongoDB Atlas",   "Cloud-hosted NoSQL database (free tier).\n"
                                "Replaces traditional SQL tables with collections"),
        ("📦  PyMongo",         "Official MongoDB driver for Python.\n"
                                "Handles all CRUD and aggregation pipelines"),
        ("🖥️  Streamlit",       "Interactive web dashboard framework.\n"
                                "7-page UI for managing the entire system"),
        ("📊  Pandas + Matplotlib", "Data manipulation and chart generation\n"
                                     "for analytics and reports"),
        ("🔐  python-dotenv",   "Securely load MongoDB credentials from .env file.\n"
                                "Prevents hardcoding passwords in source code"),
    ]
    for i, (name, desc) in enumerate(techs):
        y = Inches(1.6 + i * 0.9)
        add_shape_rect(slide, Inches(0.8), y, Inches(8.2), Inches(0.75), CARD_BG)
        add_textbox(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.65),
                    name, font_size=15, bold=True, color=ACCENT_GREEN)
        add_textbox(slide, Inches(3.6), y + Inches(0.05), Inches(5.2), Inches(0.65),
                    desc, font_size=12, color=LIGHT_GRAY)

    add_slide_number(slide, 3)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 4: PROJECT ARCHITECTURE
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Project Architecture", font_size=32, bold=True, color=ACCENT_ORANGE)
    add_accent_bar(slide, Inches(1.2), ACCENT_ORANGE)

    # Architecture layers
    layers = [
        ("🖥️  DASHBOARD LAYER", "dashboard/app.py", "Streamlit UI  –  7 interactive pages", ACCENT_BLUE),
        ("⚙️  SERVICE LAYER", "services/", "Triggers + Stored Procedures\n(order_service, alert_service, validation_service, turnaround_service)", ACCENT_GREEN),
        ("📊  QUERY LAYER", "queries/", "Aggregation Pipelines replacing SQL\n(analytics_queries, report_queries, alert_queries)", ACCENT_ORANGE),
        ("📦  MODEL LAYER", "models/", "Collection schemas + CRUD operations\n(patient_model, lab_test_model, order_model)", ACCENT_PURPLE),
        ("🗄️  DATABASE LAYER", "db/mongo_connection.py", "MongoDB Atlas connection + JSON Schema Validators", RGBColor(0xE0, 0x40, 0x40)),
    ]
    for i, (name, files, desc, color) in enumerate(layers):
        y = Inches(1.6 + i * 1.1)
        add_shape_rect(slide, Inches(0.8), y, Inches(8.2), Inches(0.95), CARD_BG, color)
        add_textbox(slide, Inches(1.0), y + Inches(0.02), Inches(3.0), Inches(0.4),
                    name, font_size=14, bold=True, color=color)
        add_textbox(slide, Inches(1.0), y + Inches(0.42), Inches(2.5), Inches(0.45),
                    files, font_size=11, color=LIGHT_GRAY)
        add_textbox(slide, Inches(4.5), y + Inches(0.1), Inches(4.3), Inches(0.75),
                    desc, font_size=12, color=WHITE)

        # Arrow between layers
        if i < len(layers) - 1:
            arrow_y = y + Inches(0.95)
            add_textbox(slide, Inches(4.5), arrow_y - Inches(0.05), Inches(1), Inches(0.25),
                        "▼", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 5: SCHEMA DESIGN (Collections)
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Schema Design – Collections", font_size=32, bold=True,
                color=ACCENT_PURPLE)
    add_accent_bar(slide, Inches(1.2), ACCENT_PURPLE)

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5),
                "4 collections with JSON Schema Validators (equivalent to CREATE TABLE with constraints)",
                font_size=14, color=LIGHT_GRAY)

    headers = ["Collection", "Key Fields", "Indexes", "SQL Constraints Applied"]
    rows = [
        ["patients", "name, age, gender,\ncontact, email, medical_history",
         "name (ASC)\ncontact (UNIQUE)", "NOT NULL, CHECK(age),\nENUM(gender), UNIQUE"],
        ["lab_tests", "test_name, category, unit,\nnormal_range, cost, TAT",
         "test_name (UNIQUE)\ncategory (ASC)", "NOT NULL, UNIQUE,\nENUM(category), CHECK(cost≥0)"],
        ["lab_orders", "patient_id, test_id, status,\npriority, result_value, dates",
         "patient_id, test_id,\nstatus, ordered_at", "FK(patient_id→patients),\nFK(test_id→lab_tests),\nENUM(status, priority)"],
        ["alerts", "order_id, alert_type,\nseverity, message, resolved",
         "order_id (ASC)\nresolved (ASC)", "FK(order_id→lab_orders),\nENUM(type, severity)"],
    ]
    add_table_slide(slide, Inches(0.5), Inches(2.2), Inches(9.0), Inches(4.0),
                    headers, rows, header_color=ACCENT_PURPLE)

    # Relationship diagram text
    add_textbox(slide, Inches(0.8), Inches(6.4), Inches(8), Inches(0.5),
                "Relationships:   patients ──< lab_orders >── lab_tests    |    lab_orders ──< alerts",
                font_size=13, bold=True, color=ACCENT_BLUE)
    add_slide_number(slide, 5)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 6: SQL TO MONGODB MAPPING
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "SQL → MongoDB Mapping", font_size=32, bold=True, color=ACCENT_BLUE)
    add_accent_bar(slide, Inches(1.2))

    headers = ["SQL Concept", "MongoDB Equivalent", "Our Implementation"]
    rows = [
        ["CREATE TABLE", "create_collection() +\nJSON Schema Validator", "mongo_connection.py\nsetup_collections()"],
        ["INSERT INTO", "insert_one({...})", "patient_model.py\ncreate_patient()"],
        ["SELECT … WHERE", "find({filter})", "All model files\nget_by_id(), search()"],
        ["UPDATE … SET", "update_one({$set: {...}})", "update_patient()\nupdate_order_status()"],
        ["JOIN … ON", "$lookup aggregation stage", "report_queries.py\nturnaround_service.py"],
        ["GROUP BY + COUNT", "$group + $sum", "analytics_queries.py\ntests_per_category()"],
        ["FOREIGN KEY", "Application-level check", "validation_service.py\nvalidate_order_data()"],
        ["BEFORE INSERT\nTrigger", "Validation function called\nbefore insert_one()", "validation_service.py"],
        ["AFTER UPDATE\nTrigger", "Function called after\nupdate_one()", "alert_service.py\ntrigger_critical_result_alert()"],
        ["Stored Procedure", "Python function wrapping\nmultiple operations", "order_service.py\nplace_order(), record_result()"],
    ]
    add_table_slide(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5),
                    headers, rows)
    add_slide_number(slide, 6)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 7: TRIGGERS & STORED PROCEDURES
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Triggers & Stored Procedures", font_size=32, bold=True,
                color=ACCENT_GREEN)
    add_accent_bar(slide, Inches(1.2), ACCENT_GREEN)

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5),
                "MongoDB has no built-in triggers — we implement them at the application level in Python",
                font_size=14, color=LIGHT_GRAY, bold=True)

    # Trigger cards
    triggers = [
        ("BEFORE INSERT Trigger", "validation_service.py",
         "• Validates patient data (name, age, gender)\n"
         "• Checks foreign key existence (patient_id, test_id)\n"
         "• Enforces status transition rules (state machine)",
         ACCENT_BLUE),
        ("AFTER UPDATE Trigger", "alert_service.py",
         "• Auto-fires when a lab result is recorded\n"
         "• Compares result against normal range\n"
         "• Creates Critical / Abnormal / Delayed alerts",
         ACCENT_ORANGE),
        ("Stored Procedure:\nsp_place_order()", "order_service.py",
         "• Validates patient & test exist (FK check)\n"
         "• Creates the order document\n"
         "• Equivalent to BEGIN TRANSACTION → INSERT → COMMIT",
         ACCENT_GREEN),
        ("Stored Procedure:\nsp_record_result()", "order_service.py",
         "• Records numeric result on order\n"
         "• Sets status to 'Completed' + timestamps\n"
         "• AFTER UPDATE trigger fires → auto-alert if abnormal",
         ACCENT_PURPLE),
    ]
    for i, (name, file, desc, color) in enumerate(triggers):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 4.7)
        top = Inches(2.2 + row * 2.5)
        add_shape_rect(slide, left, top, Inches(4.4), Inches(2.2), CARD_BG, color)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                    Inches(4.0), Inches(0.5),
                    name, font_size=14, bold=True, color=color)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.55),
                    Inches(4.0), Inches(0.3),
                    f"📁 {file}", font_size=11, color=ACCENT_BLUE)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.85),
                    Inches(4.0), Inches(1.2),
                    desc, font_size=11, color=LIGHT_GRAY)

    add_slide_number(slide, 7)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 8: AGGREGATION QUERIES
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Aggregation Queries (SQL Replacements)", font_size=32,
                bold=True, color=ACCENT_ORANGE)
    add_accent_bar(slide, Inches(1.2), ACCENT_ORANGE)

    headers = ["Query Function", "SQL Equivalent", "MongoDB Stages Used"]
    rows = [
        ["tests_per_category()", "SELECT category, COUNT(*)\nFROM lab_tests GROUP BY category",
         "$group → $project → $sort"],
        ["revenue_by_test()", "SELECT test_name, SUM(cost)\n… JOIN … GROUP BY test_name",
         "$match → $lookup → $unwind\n→ $group → $sort"],
        ["monthly_order_trends()", "SELECT DATE_FORMAT(…, '%Y-%m'),\nCOUNT(*) … GROUP BY month",
         "$group (with $year, $month)\n→ $project ($concat) → $sort"],
        ["patient_full_report()", "SELECT p.*, o.*, t.*\n… JOIN … JOIN … WHERE p.id=?",
         "$match → $lookup × 2\n→ $unwind → $group"],
        ["alert_resolution_rate()", "SELECT COUNT(*),\nSUM(CASE WHEN resolved…)",
         "$group → $project\n($cond, $divide, $multiply)"],
        ["get_turnaround_report()", "SELECT AVG(TIMESTAMPDIFF(…))\n… JOIN … GROUP BY test",
         "$match → $addFields ($subtract)\n→ $lookup → $group ($avg)"],
    ]
    add_table_slide(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.2),
                    headers, rows, header_color=ACCENT_ORANGE)

    add_textbox(slide, Inches(0.8), Inches(6.9), Inches(8), Inches(0.4),
                "Total: 12+ unique aggregation pipelines  •  10+ MongoDB stages used",
                font_size=13, color=ACCENT_GREEN, bold=True)
    add_slide_number(slide, 8)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 9: STREAMLIT DASHBOARD
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Streamlit Dashboard – 7 Pages", font_size=32, bold=True,
                color=ACCENT_PURPLE)
    add_accent_bar(slide, Inches(1.2), ACCENT_PURPLE)

    pages = [
        ("📊 Dashboard", "KPI metric cards, order/test charts,\nactive alert warnings", ACCENT_BLUE),
        ("👤 Patients", "Full CRUD: add, view, search,\ndelete patient records", ACCENT_GREEN),
        ("🔬 Lab Tests", "Define test catalog with normal\nranges, costs, categories", ACCENT_ORANGE),
        ("📋 Orders", "Place orders, update status (state\nmachine), record lab results", ACCENT_PURPLE),
        ("🔔 Alerts", "View unresolved alerts, resolve them,\ncheck for delayed orders", RGBColor(0xE0, 0x40, 0x40)),
        ("📈 Analytics", "Revenue charts, priority breakdown,\nmonthly trends, turnaround times", ACCENT_BLUE),
        ("📄 Reports", "Patient full report (multi-JOIN),\ntest result summary statistics", ACCENT_GREEN),
    ]
    for i, (name, desc, color) in enumerate(pages):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 3.1)
        top = Inches(1.6 + row * 2.3)
        w, h = Inches(2.8), Inches(2.0)
        add_shape_rect(slide, left, top, w, h, CARD_BG, color)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.15),
                    Inches(2.6), Inches(0.5),
                    name, font_size=16, bold=True, color=color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.7),
                    Inches(2.6), Inches(1.0),
                    desc, font_size=12, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # 7th card placed in the center of the last row
    add_textbox(slide, Inches(0.8), Inches(6.3), Inches(8), Inches(0.5),
                "Run with:   streamlit run dashboard/app.py",
                font_size=14, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 9)

    # ──────────────────────────────────────────────────────────────────
    #  SLIDE 10: CONCLUSION & SUMMARY
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, SECTION_BG)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                "Conclusion & Key Takeaways", font_size=32, bold=True,
                color=WHITE)
    add_accent_bar(slide, Inches(1.2), ACCENT_GREEN)

    takeaways = [
        "Successfully mapped all relational DBMS concepts to MongoDB (NoSQL)",
        "Implemented Triggers at the application level using Python service functions",
        "Created Stored Procedures that wrap validation + CRUD + trigger logic",
        "Built 12+ Aggregation Pipelines replacing SQL GROUP BY, JOIN, HAVING queries",
        "Enforced data integrity using JSON Schema Validators (NOT NULL, CHECK, ENUM, UNIQUE)",
        "Foreign Key constraints validated at the application level (validation_service.py)",
        "Full Streamlit dashboard with 7 interactive pages for real-time management",
        "Auto-alert system fires on critical / abnormal lab results and delayed orders",
    ]
    add_bullet_textbox(slide, Inches(0.8), Inches(1.6), Inches(8.2), Inches(3.5),
                       takeaways, font_size=15, color=WHITE, bullet_color=ACCENT_GREEN,
                       spacing=Pt(10))

    # Stats summary
    stats = [
        ("4", "Collections"),
        ("12+", "Aggregation\nPipelines"),
        ("4", "Triggers"),
        ("3", "Stored\nProcedures"),
    ]
    for i, (num, label) in enumerate(stats):
        left = Inches(0.8 + i * 2.25)
        add_shape_rect(slide, left, Inches(5.5), Inches(2.0), Inches(1.3),
                       CARD_BG, ACCENT_GREEN)
        add_textbox(slide, left + Inches(0.1), Inches(5.6), Inches(1.8), Inches(0.6),
                    num, font_size=30, bold=True, color=ACCENT_GREEN,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), Inches(6.15), Inches(1.8), Inches(0.5),
                    label, font_size=13, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(7.0), Inches(7), Inches(0.4),
                "Thank You!  🧪  Questions?",
                font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 10)

    # ──────────────────────────────────────────────────────────────────
    #  SAVE
    # ──────────────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__),
                               "Lab_Test_Management_Presentation.pptx")
    prs.save(output_path)
    print(f"\n✅  Presentation saved to:\n   {output_path}\n")
    return output_path


if __name__ == "__main__":
    build_presentation()
